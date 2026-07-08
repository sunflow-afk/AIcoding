#!/usr/bin/env python3
"""Generate AutoCare AI PowerPoint presentation from HTML analysis page."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ========== CONSTANTS ==========
OUTPUT = os.path.join(os.path.dirname(__file__), "AutoCareAI_前期分析_v3.0.pptx")

# Colors (from HTML dark theme)
BG = RGBColor(0x0A, 0x0E, 0x14)
SURFACE = RGBColor(0x12, 0x16, 0x1E)
CARD_BG = RGBColor(0x18, 0x1C, 0x26)
BORDER = RGBColor(0x25, 0x2B, 0x36)
TEXT = RGBColor(0xC8, 0xCC, 0xD4)
HEAD = RGBColor(0xE8, 0xEA, 0xEF)
MUTED = RGBColor(0x6B, 0x73, 0x85)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
GREEN_DIM = RGBColor(0x1A, 0x7A, 0x42)
AMBER = RGBColor(0xF0, 0xA8, 0x30)
RED = RGBColor(0xE7, 0x4C, 0x3C)
BLUE = RGBColor(0x34, 0x98, 0xDB)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ========== HELPERS ==========

def add_bg(slide, color=BG):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=12, color=TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box and return (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf

def add_para(tf, text, font_size=12, color=TEXT, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Microsoft YaHei", space_before=0, space_after=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER, border_width=1):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_slide_number(slide, num, total=8):
    """Add slide number at bottom right."""
    add_textbox(slide, 11.8, 7.0, 1.3, 0.4,
                f"{num} / {total}", font_size=9, color=MUTED, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, label, title, subtitle=None):
    """Add standard section header."""
    add_textbox(slide, 0.8, 0.4, 11.7, 0.35, label, font_size=10, color=GREEN, bold=True)
    add_textbox(slide, 0.8, 0.75, 11.7, 0.6, title, font_size=26, color=HEAD, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 1.25, 8, 0.5, subtitle, font_size=12, color=MUTED)

def add_top_bar(slide, color=GREEN):
    """Add a thin colored line at the top of the slide."""
    shape = add_rect(slide, 0, 0, 13.333, 0.06, fill_color=color, border_color=None)

def add_bottom_line(slide):
    """Add a subtle bottom separator."""
    add_rect(slide, 0, 7.35, 13.333, 0.02, fill_color=BORDER, border_color=None)


# ========== SLIDE 1: 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_top_bar(slide)

# Decorative gradient shapes
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=BG, border_color=None)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x0F, 0x1A, 0x14)
circle.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(4), Inches(4), Inches(4))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x10, 0x18, 0x22)
circle2.line.fill.background()

# Badge
badge = add_rounded_rect(slide, 0.8, 1.0, 3.0, 0.45, fill_color=CARD_BG)
badge_tf = badge.text_frame
badge_tf.word_wrap = True
p = badge_tf.paragraphs[0]
p.text = "● 尼日利亚 · 创新大赛 · 2026"
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# Title
add_textbox(slide, 0.8, 1.8, 7.5, 1.6,
            "让每一辆二手车\n都值得信赖", font_size=44, color=WHITE, bold=True)

# Subtitle
add_textbox(slide, 0.8, 3.6, 6.5, 1.0,
            "AutoCare AI 面向尼日利亚 1,400 万车主，融合 AI 诊断、维修师匹配、"
            "车辆档案与保养提醒，从问题识别到长期维护，构建二手车信任基础设施。",
            font_size=14, color=MUTED)

# Stats row
stats = [
    ("1,400 万", "汽车保有量", GREEN),
    ("75–95%", "二手车占比", HEAD),
    ("45–65%", "传音手机覆盖", HEAD),
    ("30%", "事故与机械故障相关", RED),
]
for i, (num, label, clr) in enumerate(stats):
    x = 0.8 + i * 2.6
    add_textbox(slide, x, 5.0, 2.3, 0.55, num, font_size=28, color=clr, bold=True)
    add_textbox(slide, x, 5.5, 2.3, 0.3, label, font_size=10, color=MUTED)

# Right side feature card
card = add_rect(slide, 8.0, 2.2, 4.5, 2.8, fill_color=CARD_BG)
tf = card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🔧  AutoCare AI 四合一平台"
p.font.size = Pt(14)
p.font.color.rgb = HEAD
p.font.bold = True
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.LEFT
add_para(tf, "", font_size=8, color=MUTED)
add_para(tf, "AI 音频诊断 → 技师智能匹配", font_size=11, color=TEXT)
add_para(tf, "车辆数字档案 → 保养主动提醒", font_size=11, color=TEXT)
add_para(tf, "", font_size=6, color=MUTED)
add_para(tf, "依托传音 Palm Store 预装", font_size=10, color=MUTED)
add_para(tf, "直达 1.7 亿月活用户", font_size=10, color=MUTED)

# Float tags
tag1 = add_rounded_rect(slide, 7.2, 1.6, 3.6, 0.4, fill_color=SURFACE)
tag1_tf = tag1.text_frame
tag1_tf.paragraphs[0].text = "📊 竞品空白：无人整合诊断+维修+档案"
tag1_tf.paragraphs[0].font.size = Pt(9)
tag1_tf.paragraphs[0].font.color.rgb = GREEN
tag1_tf.paragraphs[0].font.name = "Microsoft YaHei"
tag1_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

tag2 = add_rounded_rect(slide, 8.5, 5.4, 3.2, 0.4, fill_color=SURFACE)
tag2_tf = tag2.text_frame
tag2_tf.paragraphs[0].text = "📱 预装渠道：Tecno/Infinix/Itel"
tag2_tf.paragraphs[0].font.size = Pt(9)
tag2_tf.paragraphs[0].font.color.rgb = BLUE
tag2_tf.paragraphs[0].font.name = "Microsoft YaHei"
tag2_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_slide_number(slide, 1)
add_bottom_line(slide)


# ========== SLIDE 2: 市场机会 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "01 · 市场机会", "为什么是尼日利亚二手车市场？",
                  "非洲最大的二手车市场，2.3 亿人口，1,400 万辆车，维修体系数十年未变——一个等待数字化的蓝海。")

cards_data = [
    ("1,400 万", "全国汽车保有量", GREEN,
     "75–95% 为进口二手车（Tokunbo），平均车龄 ~14 年。仅 2% 人口能负担新车，年需求 72 万辆 vs 本地产量 1.4 万辆。",
     "来源：NBS, FRSC NVIS, PwC Nigeria, NADDC"),
    ("71%+", "非正规路边维修", RED,
     "无认证、无标准、无记录。95%+ 维修由路边作坊完成，60%+ 配件不合格，同一故障报价可差 10 倍。",
     "来源：Mecho Autotech, SON, FRSC, IndependentNG"),
    ("65%", "传音手机份额", BLUE,
     "传音（Tecno/Infinix/Itel）占尼日利亚智能手机 45–65%。Palm Store 月活 1.7 亿，PalmPay 3,500 万用户。",
     "来源：Canalys Africa Smartphone Pulse, Transsion Holdings"),
    ("₦900+/升", "油价飙升驱动修旧如旧", AMBER,
     "油价从 2015 年 ₦87/升飙升至 2024 年 ₦900–1,200/升（+900%）。车主更倾向修车而非换车——维修需求刚性、高频、不可推迟。",
     "来源：NBS PMS Price Watch, NMDPRA"),
]

for i, (num, label, clr, desc, src) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.1
    y = 2.0 + row * 2.55

    card = add_rect(slide, x, y, 5.8, 2.3, fill_color=CARD_BG)
    # Top accent line
    add_rect(slide, x, y, 5.8, 0.05, fill_color=clr, border_color=None)

    # Number
    add_textbox(slide, x + 0.3, y + 0.2, 5.2, 0.65, num, font_size=30, color=clr, bold=True)
    # Label
    add_textbox(slide, x + 0.3, y + 0.8, 5.2, 0.35, label, font_size=13, color=HEAD, bold=True)
    # Description
    add_textbox(slide, x + 0.3, y + 1.2, 5.2, 0.65, desc, font_size=10, color=MUTED)
    # Source
    add_textbox(slide, x + 0.3, y + 1.9, 5.2, 0.25, src, font_size=7, color=MUTED)

add_slide_number(slide, 2)
add_bottom_line(slide)


# ========== SLIDE 3: 市场数据图表 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "02 · 市场数据", "关键趋势：车越老、油越贵、进口越难",
                  "三个结构性趋势叠加，修车需求持续放大——这是 AutoCare AI 的市场基础。")

# Chart 1: Vehicle Parc & Age (top-left)
chart_data = CategoryChartData()
chart_data.categories = ['2010', '2015', '2018', '2020', '2022', '2024']
chart_data.add_series('保有量(百万辆)', (8, 11.3, 12, 11.8, 13, 14))
chart_data.add_series('平均车龄(年)', (10, 12, 13, 13, 14, 14))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.4),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = 2  # bottom
chart.legend.font.size = Pt(8)
chart.legend.font.color.rgb = MUTED

# Style chart 1
plot = chart.plots[0]
series0 = plot.series[0]
series0.format.fill.solid()
series0.format.fill.fore_color.rgb = GREEN
series1 = plot.series[1]
series1.format.fill.solid()
series1.format.fill.fore_color.rgb = RED

chart.chart_title.has_text_frame = True
chart.chart_title.text_frame.paragraphs[0].text = "汽车保有量与平均车龄变迁"
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart.chart_title.text_frame.paragraphs[0].font.color.rgb = HEAD

# Dark background for chart area
chart_frame.chart.font.color.rgb = MUTED
chart_frame.chart.font.size = Pt(8)

add_textbox(slide, 6.65, 4.3, 5.5, 0.2,
            "来源：NBS Road Transport Data, FRSC NVIS, NADDC", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

# Chart 2: Fuel Price (top-right)
chart_data2 = CategoryChartData()
chart_data2.categories = ['2015','2016','2017','2018','2019','2020','2021','2022','2023Q1','2023Q2','2023Q4','2024']
chart_data2.add_series('油价 (₦/升)', (87, 145, 146, 146, 146, 167, 170, 206, 195, 557, 672, 950))

chart_frame2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.4),
    chart_data2
)
chart2 = chart_frame2.chart
chart2.has_legend = True
chart2.legend.position = 2
chart2.legend.font.size = Pt(8)
chart2.legend.font.color.rgb = MUTED

series = chart2.series[0]
series.format.line.color.rgb = AMBER
series.format.line.width = Pt(2)

chart2.chart_title.has_text_frame = True
chart2.chart_title.text_frame.paragraphs[0].text = "燃油价格飙升（₦/升）"
chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart2.chart_title.text_frame.paragraphs[0].font.color.rgb = HEAD

chart_frame2.chart.font.color.rgb = MUTED
chart_frame2.chart.font.size = Pt(8)

add_textbox(slide, 12.85, 4.3, 5.5, 0.2,
            "来源：NBS PMS Price Watch 月度报告", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

# Chart 3: Import Value vs Volume (bottom-left)
chart_data3 = CategoryChartData()
chart_data3.categories = ['2017','2018','2019','2020','2021','2022','2023','2024','2025']
chart_data3.add_series('进口额(₦万亿)', (0.145, 0.269, 0.58, 0.65, 0.52, 0.31, 1.47, 1.26, 1.58))

chart_frame3 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(4.6), Inches(5.8), Inches(2.4),
    chart_data3
)
chart3 = chart_frame3.chart
chart3.has_legend = True
chart3.legend.position = 2
chart3.legend.font.size = Pt(8)
chart3.legend.font.color.rgb = MUTED

series3 = chart3.series[0]
series3.format.fill.solid()
series3.format.fill.fore_color.rgb = BLUE

chart3.chart_title.has_text_frame = True
chart3.chart_title.text_frame.paragraphs[0].text = "二手车进口额：量缩价涨"
chart3.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart3.chart_title.text_frame.paragraphs[0].font.color.rgb = HEAD

chart_frame3.chart.font.color.rgb = MUTED
chart_frame3.chart.font.size = Pt(8)

add_textbox(slide, 6.65, 6.9, 5.5, 0.2,
            "来源：NBS Foreign Trade, NPA Shipping, US CBP", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

# Chart 4: Network Stacked (bottom-right)
chart_data4 = CategoryChartData()
chart_data4.categories = ['2017', '2019', '2021', '2023', '2025']
chart_data4.add_series('2G', (72, 60, 50, 45, 40))
chart_data4.add_series('3G', (18, 25, 22, 18, 7))
chart_data4.add_series('4G', (10, 15, 28, 35, 50))
chart_data4.add_series('5G', (0, 0, 0, 2, 3))

chart_frame4 = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED,
    Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.4),
    chart_data4
)
chart4 = chart_frame4.chart
chart4.has_legend = True
chart4.legend.position = 2
chart4.legend.font.size = Pt(8)
chart4.legend.font.color.rgb = MUTED

colors4 = [MUTED, PURPLE, BLUE, GREEN]
for i, s in enumerate(chart4.series):
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = colors4[i]

chart4.chart_title.has_text_frame = True
chart4.chart_title.text_frame.paragraphs[0].text = "网络连接结构变迁（%）"
chart4.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart4.chart_title.text_frame.paragraphs[0].font.color.rgb = HEAD

chart_frame4.chart.font.color.rgb = MUTED
chart_frame4.chart.font.size = Pt(8)

add_textbox(slide, 12.85, 6.9, 5.5, 0.2,
            "来源：NCC, GSMA Mobile Economy SSA 2024", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

add_slide_number(slide, 3)
add_bottom_line(slide)


# ========== SLIDE 4: 6大痛点 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "03 · 用户痛点", "六大核心痛点，系统性的信任真空",
                  "调研覆盖 Lagos 与 Abuja 车主。弱网限制和低付费意愿已归入设计约束，另作讨论。")

pains = [
    ("痛点 01", "★★★★★", "🔍 车辆问题难判断",
     "技师用\"试错法\"诊断，水泵故障 ₦45k 被报\"发动机大修\"₦710k。7/10 车主反映被反复误诊。高频：每月 1–3 次", RED),
    ("痛点 02", "★★★★★", "👨‍🔧 维修师难选择",
     "催化器价值 ₦100 万被偷换成 ₦2.5 万假货。95%+ 维修由非正规技师完成，无认证。全靠熟人推荐——但熟人也不一定懂", RED),
    ("痛点 03", "★★★★★", "💰 报价不透明",
     "同一故障报价差 10 倍，₦600 垫片报 ₦22,000。女性车主报价系统性偏高。技师与配件商勾结，无书面报价单", RED),
    ("痛点 04", "★★★★☆", "📋 维修记录缺失",
     "修了什么、换了什么——全无记录。\"信息黑洞\"延续数十年，事故车可轻易流入二手车市场。卖车时买家无从验证车况", AMBER),
    ("痛点 05", "★★★★☆", "📅 保养提醒不足",
     "多数车主\"开到坏了才修\"，小问题拖成大故障。~60% 车主连轮胎气压表都不信任。技师不提醒——修大故障更赚钱", AMBER),
    ("痛点 06", "★★★★★", "🔧 假配件泛滥",
     "60%+ 配件不合格。翻新件冒充新件，过期轮胎重新喷漆。正品价三年涨 150–300%。30% 致命事故源于机械故障（FRSC）", AMBER),
]

for i, (num, stars, title, desc, accent) in enumerate(pains):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.9 + row * 2.45

    card = add_rect(slide, x, y, 3.8, 2.15, fill_color=CARD_BG)
    add_rect(slide, x, y, 3.8, 0.04, fill_color=accent, border_color=None)

    add_textbox(slide, x + 0.25, y + 0.15, 3.3, 0.25, num, font_size=8, color=MUTED, bold=True)
    add_textbox(slide, x + 0.25, y + 0.35, 3.3, 0.2, stars, font_size=9, color=AMBER)
    add_textbox(slide, x + 0.25, y + 0.55, 3.3, 0.3, title, font_size=12, color=HEAD, bold=True)
    add_textbox(slide, x + 0.25, y + 0.95, 3.3, 1.0, desc, font_size=9, color=MUTED)

# Constraint note
note = add_rect(slide, 0.8, 6.8, 11.7, 0.45, fill_color=SURFACE, border_color=AMBER)
note_tf = note.text_frame
note_tf.word_wrap = True
p = note_tf.paragraphs[0]
p.text = ("📐 设计约束（非痛点）：弱网/低端设备限制（2G 仍占 40%+，1GB=₦431）和低付费意愿"
          "（月最低工资 ₦70,000）——将在技术方案和商业模式中解决")
p.font.size = Pt(9)
p.font.color.rgb = MUTED
p.font.name = "Microsoft YaHei"

add_slide_number(slide, 4)
add_bottom_line(slide)


# ========== SLIDE 5: 用户画像 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "04 · 用户画像", "两个核心画像，覆盖最广泛用车场景",
                  "Emeka（生存驱动型，车 = 饭碗）+ Chidi（信息弱势型，怕被宰、不懂车）")

personas = [
    {
        "name": "Emeka Okonkwo",
        "emoji": "👨🏾‍🦱",
        "role": "Lagos · 网约车司机 · 34 岁",
        "quote": "\"车是我的饭碗。它不跑，我就没饭吃。我不相信任何技师——但我没有选择。\"",
        "fields": [
            ("家庭", "已婚，2 个孩子（6 岁、3 岁），妻子在市场卖布料"),
            ("月收入", "毛收入 ₦300k–400k → 净收入仅 ₦100k–180k"),
            ("车辆", "2012 款丰田卡罗拉（Tokunbo），日均 200–250km"),
            ("手机", "Tecno Spark 20 · MTN 4G · 3GB/月"),
            ("工作节奏", "6:00–21:00，14–15h/天，Bolt+Uber+InDrive 三平台"),
            ("核心需求", "🚨 故障分诊 👨‍🔧 靠谱技师 💰 费用透明 📊 收支管理"),
        ],
        "insight": "💡 用户洞察：Emeka 不需要一个工具，他需要一个站在他那边的\"汽车顾问\"——"
                   "告诉他车到底怎么了、该花多少钱、该找谁修。生存驱动，高频刚需。",
        "evidence": [
            "Bolt/Ipsos 2024：尼日利亚 72 万网约车司机，64% 报告生活质量改善",
            "真实案例 — Noel Onoja：被 mechanic 反复误诊——气流计 ₦5k→喷油嘴 ₦6k→\"换大脑板\"",
            "Malomo (2024) Lagos 调研：79% 男性，50% 年龄 22–32 岁",
        ],
        "accent": GREEN,
        "avatar_bg": GREEN,
    },
    {
        "name": "Chidi Eze",
        "emoji": "👨🏾‍💼",
        "role": "Abuja · 日常私家车主 · 38 岁",
        "quote": "\"每次修车我都不知道是不是被宰了。我需要一个告诉我'正常价是多少'的人。\"",
        "fields": [
            ("家庭", "已婚，3 个孩子，妻子是小学教师"),
            ("月收入", "家庭月总收入约 ₦380k（他 ₦250k + 妻子 ₦130k）"),
            ("车辆", "2010 款丰田凯美瑞（Tokunbo），里程 18 万 km，周末兼跑 Bolt"),
            ("手机", "Infinix Hot 40 · Airtel 网络 · 2.5GB/月"),
            ("修车经历", "换刹车片被收 ₦85k，后来才知道正常价 ₦30k–45k"),
            ("核心需求", "💰 报价参考 🔍 故障预判 📅 保养提醒 📋 维修记录"),
        ],
        "insight": "💡 用户洞察：Chidi 懂一点车——但正因为\"懂一点\"，他更清楚自己有多容易被骗。"
                   "他需要一个随时能掏出来的\"底价参考\"，让他在技师面前有底气说\"我查过，正常价不是这样\"。",
        "evidence": [
            "真实案例 — Star Okigwe：攒 18 个月买车，3 个月即抛锚——\"Tokunbo 车从买来第一天就要修\"",
            "Sunday Balogun：\"我让技师告诉我需要换什么，自己去市场买配件再让他装\"",
            "Edo State 调研（250 名车主）：56% 维修为机械故障，满意度仅\"中等偏上\"",
        ],
        "accent": BLUE,
        "avatar_bg": BLUE,
    },
]

for i, p in enumerate(personas):
    x = 0.8 + i * 6.1

    # Card background
    add_rect(slide, x, 1.8, 5.8, 5.4, fill_color=CARD_BG)

    # Header
    header = add_rect(slide, x, 1.8, 5.8, 1.8, fill_color=SURFACE)
    # Avatar circle
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 2.15), Inches(2.0), Inches(0.7), Inches(0.7))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = p["avatar_bg"]
    avatar.line.fill.background()
    avatar_tf = avatar.text_frame
    avatar_tf.paragraphs[0].text = p["emoji"]
    avatar_tf.paragraphs[0].font.size = Pt(20)
    avatar_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + 0.3, 2.8, 5.2, 0.3, p["name"], font_size=15, color=HEAD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.1, 5.2, 0.25, p["role"], font_size=10, color=p["accent"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.4, 3.35, 5.0, 0.35, p["quote"], font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Fields (2 columns)
    for j, (label, value) in enumerate(p["fields"]):
        col = j % 2
        row = j // 2
        fx = x + 0.25 + col * 2.85
        fy = 3.75 + row * 0.5
        add_textbox(slide, fx, fy, 2.6, 0.15, label, font_size=7, color=MUTED, bold=True)
        add_textbox(slide, fx, fy + 0.16, 2.6, 0.3, value, font_size=9, color=HEAD)

    # Insight
    insight_y = 5.35
    insight_bg = add_rect(slide, x + 0.3, insight_y, 5.2, 0.65, fill_color=RGBColor(0x0F, 0x24, 0x18), border_color=GREEN_DIM)
    add_textbox(slide, x + 0.5, insight_y + 0.05, 4.8, 0.55, p["insight"], font_size=8, color=GREEN)

    # Evidence
    evid_y = 6.1
    evid_bg = add_rect(slide, x + 0.3, evid_y, 5.2, 1.0, fill_color=RGBColor(0x24, 0x1A, 0x0F), border_color=AMBER)
    add_textbox(slide, x + 0.5, evid_y + 0.05, 4.8, 0.15, "📰 真实证据", font_size=8, color=AMBER, bold=True)
    for k, ev in enumerate(p["evidence"]):
        add_textbox(slide, x + 0.5, evid_y + 0.25 + k * 0.22, 4.8, 0.2, f"• {ev}", font_size=7, color=MUTED)

add_slide_number(slide, 5)
add_bottom_line(slide)


# ========== SLIDE 6: 旅程对比 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "05 · 用户旅程对比", "一条旅程，两个世界",
                  "融合真实报道事件——左列\"现状\"，右列\"使用 AutoCare AI 后\"。最大情绪差出现在故障判断和长期维护阶段。")

journey = [
    ("❶ 日常运营", "层层抽成，勉强维生\n油价 ₦900+/升，平台抽 25%，agbero levy ₦1,500/天",
     "数据驱动，心中有数\nAutoCare AI 后台追踪里程+油耗，月末自动生成收支报告"),
    ("❷ 发现异常", "完全看不懂，极度焦虑\n靠边打开引擎盖——完全看不懂。\"还能接单吗？\"",
     "2 分钟，3MB，搞定\n录制 30 秒音频 + 拍仪表盘照片，2 分钟完成诊断"),
    ("❸ 故障判断 📰", "猜谜式诊断，花冤枉钱\nNoel Onoja：换气流计 ₦5k→无效→换喷油嘴 ₦6k→仍无效",
     "15 秒出结果，⚡ 核心价值\n🟡 85% 概率水泵轴承磨损，预估 ₦35k–55k，非紧急"),
    ("❹ 找技师", "熟人推荐 = 随机赌博\nWhatsApp 群问\"谁认识靠谱 mechanic？\"——等于赌博",
     "四维匹配，一键预约\n推荐 3 家评分 4.5+，按距离+评分+专长+估价排序"),
    ("❺ 到店维修 📰", "无价格参照，无法质疑\nSunday Balogun：\"我自己买配件再让技师装——再也不信技师代买\"",
     "双重信任，安心交车\n技师确认：\"You're right, na water pump.\" 报价在预估范围内 ✅"),
    ("❻ 记录与评价", "零记录，\"信息黑洞\"\nTimothy Raphael：修完无记录，三月后同问题复发",
     "一键存档，数字档案\n拍照 OCR 自动提取 → 维修记录存入车辆数字档案，卖车 = 溢价凭证"),
    ("❼ 长期维护", "开坏了才修，没人关心你的车\n~60% 车主连气压表都不信任。30% 致命事故与机械故障有关",
     "从\"修车工具\"到\"养车管家\"\n3 个月后推送保养提醒→一键预约。\"比我自己还清楚我的车。\" 🔁"),
]

# Column headers
header_red = add_rect(slide, 0.8, 1.8, 5.6, 0.4, fill_color=RGBColor(0x2A, 0x10, 0x10))
header_red_tf = header_red.text_frame
header_red_tf.paragraphs[0].text = "🔴 现状（无 AutoCare AI）"
header_red_tf.paragraphs[0].font.size = Pt(11)
header_red_tf.paragraphs[0].font.color.rgb = RED
header_red_tf.paragraphs[0].font.name = "Microsoft YaHei"
header_red_tf.paragraphs[0].bold = True
header_red_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

header_green = add_rect(slide, 6.9, 1.8, 5.6, 0.4, fill_color=RGBColor(0x10, 0x2A, 0x15))
header_green_tf = header_green.text_frame
header_green_tf.paragraphs[0].text = "🟢 使用 AutoCare AI 后"
header_green_tf.paragraphs[0].font.size = Pt(11)
header_green_tf.paragraphs[0].font.color.rgb = GREEN
header_green_tf.paragraphs[0].font.name = "Microsoft YaHei"
header_green_tf.paragraphs[0].bold = True
header_green_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

for i, (stage, current, ideal) in enumerate(journey):
    y = 2.3 + i * 0.7

    # Current column
    cell_cur = add_rect(slide, 0.8, y, 5.6, 0.65, fill_color=CARD_BG)
    add_textbox(slide, 0.95, y + 0.02, 5.3, 0.18, stage, font_size=9, color=RED, bold=True)
    add_textbox(slide, 0.95, y + 0.22, 5.3, 0.38, current, font_size=8, color=MUTED)

    # Ideal column
    cell_ideal = add_rect(slide, 6.9, y, 5.6, 0.65, fill_color=CARD_BG)
    add_textbox(slide, 7.05, y + 0.02, 5.3, 0.18, stage, font_size=9, color=GREEN, bold=True)
    add_textbox(slide, 7.05, y + 0.22, 5.3, 0.38, ideal, font_size=8, color=MUTED)

# Emotion comparison table
emotion_data = [
    ("🔴 现状", ["😊→😤", "😰", "😣→😡", "😟", "😣😤", "😞", "😞→😰"]),
    ("🟢 AutoCare AI", ["😊→😤", "😐→🤔", "😌💪", "🙂👍", "😊🙏", "😄📝", "😊🔁💡"]),
    ("📊 情绪差", ["起点相同", "+2", "+5 🔑", "+2", "+3", "+3", "+5 🔑"]),
]

stages_labels = ["❶ 日常运营", "❷ 发现异常", "❸ 故障判断", "❹ 找技师", "❺ 到店维修", "❻ 记录评价", "❼ 长期维护"]

# Table header
th_y = 7.2
for j, lbl in enumerate(stages_labels):
    x = 1.3 + j * 1.6
    add_textbox(slide, x, th_y, 1.5, 0.22, lbl, font_size=5.5, color=HEAD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, th_y, 0.5, 0.22, "", font_size=5, color=MUTED)

add_bottom_line(slide)
add_slide_number(slide, 6)


# ========== SLIDE 7: 竞争格局 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "06 · 竞争格局", "没有直接竞品——但窗口不会永远敞开",
                  "现有玩家各做一段，无人整合\"诊断 + 匹配 + 档案 + 提醒\"全链条。传音渠道 + 弱网优先 = 为尼日利亚而建。")

# Competition table
comp_data = [
    ("在线交易", "Jiji Cars, Cars45", ["✕", "✕", "✕", "✕", "✕", "✕"]),
    ("配件商城", "Jiji, Ladipo 线上", ["✕", "✕", "✕", "✕", "✕", "✕"]),
    ("技师目录", "Carlots.ng", ["✕", "✓", "✕", "✕", "✕", "✕"]),
    ("通用 AI", "ChatGPT, Gemini", ["⚠ 泛化", "✕", "✕", "✕", "✕", "✕"]),
    ("传统线下", "Ladipo/Apo 集群", ["✕", "⚠ 熟人", "✕", "✕", "✕", "✕"]),
    ("叫车平台", "Uber/Bolt 内置", ["✕", "✕", "✕", "✕", "✕", "✕"]),
]

headers = ["类型", "代表产品", "AI 诊断", "技师匹配", "车辆档案", "保养提醒", "弱网优化", "传音预装"]
col_widths = [1.2, 2.0, 1.15, 1.15, 1.15, 1.15, 1.15, 1.15]

# Table header row
th_x = 0.8
for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    th = add_rect(slide, th_x, 2.0, w, 0.4, fill_color=SURFACE)
    th_tf = th.text_frame
    th_tf.paragraphs[0].text = hdr
    th_tf.paragraphs[0].font.size = Pt(10)
    th_tf.paragraphs[0].font.color.rgb = HEAD
    th_tf.paragraphs[0].font.bold = True
    th_tf.paragraphs[0].font.name = "Microsoft YaHei"
    th_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    th_x += w

# Table data rows
for i, (type_name, name, checks) in enumerate(comp_data):
    y = 2.5 + i * 0.45
    row_data = [type_name, name] + checks
    td_x = 0.8
    for j, (val, w) in enumerate(zip(row_data, col_widths)):
        td = add_rect(slide, td_x, y, w, 0.4, fill_color=CARD_BG if i % 2 == 0 else SURFACE)
        td_tf = td.text_frame
        td_tf.paragraphs[0].text = val
        td_tf.paragraphs[0].font.size = Pt(9)
        td_tf.paragraphs[0].font.name = "Microsoft YaHei"
        td_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if val == "✓":
            td_tf.paragraphs[0].font.color.rgb = GREEN
            td_tf.paragraphs[0].font.bold = True
        elif val == "✕":
            td_tf.paragraphs[0].font.color.rgb = MUTED
        else:
            td_tf.paragraphs[0].font.color.rgb = TEXT
        td_x += w

# AutoCare AI row (highlighted)
us_y = 2.5 + 6 * 0.45
us_values = ["AutoCare AI", "本项目 · 2026", "✓", "✓", "✓", "✓", "✓", "✓"]
td_x = 0.8
for j, (val, w) in enumerate(zip(us_values, col_widths)):
    td = add_rect(slide, td_x, us_y, w, 0.45, fill_color=RGBColor(0x10, 0x24, 0x19))
    td_tf = td.text_frame
    td_tf.paragraphs[0].text = val
    td_tf.paragraphs[0].font.size = Pt(10)
    td_tf.paragraphs[0].font.name = "Microsoft YaHei"
    td_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if val == "✓":
        td_tf.paragraphs[0].font.color.rgb = GREEN
        td_tf.paragraphs[0].font.bold = True
    elif j == 0:
        td_tf.paragraphs[0].font.color.rgb = HEAD
        td_tf.paragraphs[0].font.bold = True
    else:
        td_tf.paragraphs[0].font.color.rgb = HEAD
    td_x += w

# Differentiation card
diff_card = add_rect(slide, 0.8, 5.6, 11.7, 1.3, fill_color=CARD_BG, border_color=GREEN)
diff_tf = diff_card.text_frame
diff_tf.word_wrap = True
p = diff_tf.paragraphs[0]
p.text = "🎯 核心差异化"
p.font.size = Pt(13)
p.font.color.rgb = GREEN
p.font.bold = True
p.font.name = "Microsoft YaHei"
add_para(diff_tf, "AI 诊断 × 技师匹配 × 车辆档案 × 保养提醒——四位一体，而非单独一个功能。", font_size=11, color=HEAD)
add_para(diff_tf, "📱 传音系统预装 + Palm Store 分发——获客成本远低于独立 App。", font_size=11, color=HEAD)
add_para(diff_tf, "🌍 弱网/低端设备优先——为尼日利亚的现实而设计，不为硅谷的 5G 用户设计。", font_size=11, color=HEAD)

add_slide_number(slide, 7)
add_bottom_line(slide)


# ========== SLIDE 8: 愿景与下一步 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_section_title(slide, "07 · 愿景与下一步", "从\"修车工具\"到\"车辆全生命周期管家\"",
                  "终极愿景：尼日利亚版\"Carfax + 大众点评修车版 + 车主管家\"——覆盖车辆全生命周期的数据平台。")

# 3 Vision cards
visions = [
    ("📋", "数字车辆档案",
     "自动累积维修记录 + 保养记录 + 事故记录 + 里程 + 费用统计。二手车交易时一键生成 AutoCare 车辆健康报告（Carfax 模式），成为买卖双方的信任凭证。"),
    ("🏦", "金融科技合作",
     "与 FairMoney（车辆融资）、UBA（$1 亿 Drive-to-Own 计划）合作，用车况数据做风控——谁掌握车辆完整数据，谁就掌握二手车交易定价权。"),
    ("🏪", "传音生态协同",
     "Palm Store 预装直达 1.7 亿 MAU → PalmPay 3,500 万用户小额支付 → Carlcare 1,200+ 服务中心作为线下诊断核验节点——线上线下打通。"),
]

for i, (icon, title, desc) in enumerate(visions):
    x = 0.8 + i * 4.1
    card = add_rect(slide, x, 2.0, 3.8, 2.2, fill_color=CARD_BG)
    add_textbox(slide, x + 0.3, 2.1, 3.2, 0.5, icon, font_size=28, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 2.6, 3.2, 0.3, title, font_size=13, color=HEAD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.0, 3.2, 1.0, desc, font_size=9, color=MUTED)

# Key assumptions table
add_textbox(slide, 0.8, 4.5, 5, 0.35, "🔑 关键假设与验证优先级", font_size=14, color=HEAD, bold=True)

assumptions = [
    ("🔴 高", "网约车司机愿意为\"减少停工\"付费", "与 Lagos AUATON 工会合作，问卷 + 深度访谈 50 名司机，A/B 定价测试"),
    ("🔴 高", "AI 音频诊断在老旧二手车上的准确率可达 80%+", "在 Lagos 与 10 家合作技师采集 500+ 真实故障音频样本"),
    ("🔴 高", "传音 Palm Store 预装可带来显著低成本获客", "与传音商务团队初步沟通合作可行性"),
    ("🟡 中", "技师愿意接受平台评价和价格透明化", "在 Lagos Ikeja/Apapa 邀请 20 位技师深度访谈"),
    ("🟡 中", "用户会在维修后主动录入维修记录", "MVP 设计拍照 OCR + 技师端确认双通道，测试 30 天留存"),
    ("🟢 低", "保养提醒可显著提升用户留存和复购", "MVP 做 A/B 测试（有提醒 vs 无提醒），对比 90 天留存率"),
]

ah_col_widths = [0.8, 4.5, 6.4]

# Table header
th_x2 = 0.8
for hdr, w in zip(["优先级", "假设", "验证方式"], ah_col_widths):
    th = add_rect(slide, th_x2, 5.0, w, 0.37, fill_color=SURFACE)
    th_tf = th.text_frame
    th_tf.paragraphs[0].text = hdr
    th_tf.paragraphs[0].font.size = Pt(10)
    th_tf.paragraphs[0].font.color.rgb = HEAD
    th_tf.paragraphs[0].font.bold = True
    th_tf.paragraphs[0].font.name = "Microsoft YaHei"
    th_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    th_x2 += w

for i, (pri, hyp, verify) in enumerate(assumptions):
    y = 5.45 + i * 0.32
    row_data = [pri, hyp, verify]
    td_x2 = 0.8
    for j, (val, w) in enumerate(zip(row_data, ah_col_widths)):
        td = add_rect(slide, td_x2, y, w, 0.3, fill_color=CARD_BG if i % 2 == 0 else SURFACE)
        td_tf = td.text_frame
        td_tf.paragraphs[0].text = val
        td_tf.paragraphs[0].font.size = Pt(8)
        td_tf.paragraphs[0].font.name = "Microsoft YaHei"
        td_tf.paragraphs[0].alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
        td_tf.paragraphs[0].font.color.rgb = RED if "🔴" in val else (AMBER if "🟡" in val else (GREEN if "🟢" in val else TEXT))
        td_x2 += w

add_slide_number(slide, 8)
add_bottom_line(slide)


# ========== FINAL SLIDE: 感谢页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)

add_textbox(slide, 0.8, 2.2, 11.7, 1.0, "谢谢！", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 3.3, 11.7, 0.6, "🚗 AutoCare AI — 尼日利亚二手车维护助手", font_size=18, color=GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 4.0, 11.7, 0.5, "前期分析 v3.0 · 数据跨度 2010–2025", font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 4.5, 11.7, 0.4,
            "完整数据来源见《AutoCareAI_前期分析.md》附录", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 5.5, 11.7, 0.3, "© 2026 AutoCare AI 团队 · 创新大赛", font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

# ========== SAVE ==========
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
